# agents/ingestion_agent.py

import os
import csv
import docx
import pptx
import markdown
from PyPDF2 import PdfReader
from mcp.protocol import MCPMessage

class IngestionAgent:
    def __init__(self, name="IngestionAgent"):
        self.name = name

    def handle(self, mcp_message: MCPMessage) -> MCPMessage:
        if mcp_message.type != "DOCUMENT_UPLOAD":
            raise ValueError(f"{self.name} cannot handle message type: {mcp_message.type}")

        file_paths = mcp_message.payload.get("file_paths", [])
        query = mcp_message.payload.get("query", "")

        all_text = []

        for path in file_paths:
            ext = os.path.splitext(path)[-1].lower()

            if ext == ".pdf":
                all_text.append(self.parse_pdf(path))
            elif ext == ".docx":
                all_text.append(self.parse_docx(path))
            elif ext == ".pptx":
                all_text.append(self.parse_pptx(path))
            elif ext == ".csv":
                all_text.append(self.parse_csv(path))
            elif ext in [".txt", ".md"]:
                all_text.append(self.parse_txt(path))
            else:
                all_text.append(f"Unsupported file format: {path}")

        return MCPMessage(
            sender=self.name,
            receiver="RetrievalAgent",
            type="DOCUMENT_PARSED",
            trace_id=mcp_message.trace_id,
            payload={
                "documents": all_text,
                "query": query
            }
        )

    def parse_pdf(self, path):
        text = ""
        with open(path, "rb") as f:
            reader = PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ""
        return text

    def parse_docx(self, path):
        doc = docx.Document(path)
        return "\n".join([p.text for p in doc.paragraphs])

    def parse_pptx(self, path):
        prs = pptx.Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def parse_csv(self, path):
        rows = []
        with open(path, newline='', encoding="utf-8") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(", ".join(row))
        return "\n".join(rows)

    def parse_txt(self, path):
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
